import PyPDF2
import docx
import csv
import io

def extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Extract text from PDF, DOCX, TXT, MD, CSV, Excel, and PPT files.
    """
    ext = filename.lower().split('.')[-1]
    
    if ext == 'pdf':
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    
    elif ext == 'docx':
        doc = docx.Document(io.BytesIO(file_bytes))
        return "\n".join([p.text for p in doc.paragraphs])
    
    elif ext in ['txt', 'md']:
        return file_bytes.decode('utf-8')
    
    elif ext == 'csv':
        text_content = file_bytes.decode('utf-8')
        reader = csv.reader(io.StringIO(text_content))
        rows = []
        for row in reader:
            rows.append(" | ".join(row))
        return "\n".join(rows)
    
    elif ext in ['xlsx', 'xls']:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
            text = ""
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text += f"--- Sheet: {sheet_name} ---\n"
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    text += " | ".join(cells) + "\n"
            wb.close()
            return text
        except ImportError:
            return "Excel support requires openpyxl package."
    
    elif ext in ['pptx']:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            text = ""
            for i, slide in enumerate(prs.slides, 1):
                text += f"--- Slide {i} ---\n"
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text += paragraph.text + "\n"
            return text
        except ImportError:
            return "PPT support requires python-pptx package."
    
    return ""

def chunk_text(text: str, chunk_size=1000, overlap=200):
    """
    Split text into overlapping chunks for indexing.
    """
    chunks = []
    start = 0
    if not text:
        return []
        
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        # Move start forward but subtract overlap to keep context
        start = end - overlap
        
        # Guard against infinite loops or tiny increments
        if start <= 0 and len(text) > chunk_size:
             start = chunk_size - overlap
        elif start >= len(text):
            break
            
    return chunks
